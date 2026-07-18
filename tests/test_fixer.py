"""Tests for the auto-fix engine."""

import os
import sys
import tempfile
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from backend.engine.fixer import Fixer, quick_fix
from backend.models.issue import Issue, ElementLocation
from backend.parsers import parse_pptx


def _make_test_pptx(path: str):
    """Create a simple test PPTX with fixable issues."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide with non-standard font
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Comic Sans text"
    p.font.name = "Comic Sans MS"
    p.font.size = Pt(14)

    # Slide with small text
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Small text body"
    p2.font.name = "微软雅黑"
    p2.font.size = Pt(8)

    prs.save(path)
    return path


def test_fixer_font_replace():
    """Test replacing a non-standard font."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        input_path = f.name
    output_path = input_path.replace('.pptx', '_fixed.pptx')

    try:
        _make_test_pptx(input_path)

        issues = [
            Issue(
                rule_id="FONT-001",
                rule_name="非标准字体",
                severity="S3",
                description="Non-standard font",
                detail="Comic Sans MS used",
                suggestion="Replace with standard",
                location=ElementLocation(
                    slide_index=0,
                    element_id=None,
                    element_name="TextBox 1",
                ),
                auto_fixable=True,
                fix_type="font_replace",
                fix_params={"old_font": "Comic Sans MS", "new_font": "微软雅黑"},
            )
        ]

        fixer = Fixer()
        fixed_count, out_path = fixer.apply_fixes(input_path, issues, output_path)

        assert fixed_count == 1, f"Expected 1 fix, got {fixed_count}"
        assert os.path.exists(out_path), "Output file not created"

        # Verify the fix
        data = parse_pptx(out_path)
        slide0 = data.slides[0]
        fonts = set()
        for s in slide0.shapes:
            for t in s.text_elements:
                if t.font_name:
                    fonts.add(t.font_name)
        assert "Comic Sans MS" not in fonts, "Font should have been replaced"
        print(f"PASS: Font replace test - fonts now: {fonts}")

    finally:
        for p in [input_path, output_path]:
            if os.path.exists(p):
                os.unlink(p)


def test_fixer_font_size():
    """Test fixing too-small text."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        input_path = f.name
    output_path = input_path.replace('.pptx', '_fixed.pptx')

    try:
        _make_test_pptx(input_path)

        issues = [
            Issue(
                rule_id="SIZE-001",
                rule_name="正文字号过小",
                severity="S3",
                description="Font size too small",
                detail="8pt < 10pt minimum",
                suggestion="Set to 10pt",
                location=ElementLocation(
                    slide_index=1,
                    element_id=None,
                    element_name="TextBox 1",
                ),
                auto_fixable=True,
                fix_type="font_size_set",
                fix_params={"target_size": 10},
            )
        ]

        fixer = Fixer()
        fixed_count, out_path = fixer.apply_fixes(input_path, issues, output_path)

        assert fixed_count == 1, f"Expected 1 fix, got {fixed_count}"

    finally:
        for p in [input_path, output_path]:
            if os.path.exists(p):
                os.unlink(p)


def test_fixer_no_fixable():
    """Test fixer with no fixable issues."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        input_path = f.name
    output_path = input_path.replace('.pptx', '_fixed.pptx')

    try:
        _make_test_pptx(input_path)

        issues = [
            Issue(
                rule_id="TEST-001",
                rule_name="Non-fixable issue",
                severity="S3",
                description="Test",
                detail="Test detail",
                suggestion="Do something",
                location=ElementLocation(slide_index=0),
                auto_fixable=False,
            )
        ]

        fixer = Fixer()
        fixed_count, out_path = fixer.apply_fixes(input_path, issues, output_path)

        assert fixed_count == 0, "Should have fixed nothing"

    finally:
        for p in [input_path, output_path]:
            if os.path.exists(p):
                os.unlink(p)


def test_quick_fix():
    """Test the quick_fix convenience function."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        input_path = f.name

    try:
        _make_test_pptx(input_path)

        issues = [
            Issue(
                rule_id="FONT-001",
                rule_name="非标准字体",
                severity="S3",
                description="Non-standard font",
                detail="Comic Sans MS used",
                suggestion="Replace",
                location=ElementLocation(
                    slide_index=0,
                    element_id=None,
                    element_name="TextBox 1",
                ),
                auto_fixable=True,
                fix_type="font_replace",
                fix_params={"old_font": "Comic Sans MS", "new_font": "微软雅黑"},
            )
        ]

        fixed_count, out_path = quick_fix(input_path, issues)
        assert fixed_count > 0, "quick_fix should have fixed something"

    finally:
        for p in [input_path]:
            if os.path.exists(p):
                os.unlink(p)
        # Clean up auto-named output
        auto_out = input_path.replace('.pptx', '_fixed.pptx')
        if os.path.exists(auto_out):
            os.unlink(auto_out)
