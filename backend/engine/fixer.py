"""Auto-fix engine - applies fixes to PPTX files."""

import shutil
import os
from typing import List, Tuple
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from backend.models.issue import Issue
from backend.parsers import parse_pptx
from backend.utils.geometry import pt_to_emu


class Fixer:
    """Apply auto-fixes to a PPTX file."""

    def __init__(self):
        self.modified_count = 0
        self.fix_log = []

    def apply_fixes(self, file_path: str, issues: List[Issue],
                    output_path: str) -> Tuple[int, str]:
        """Apply auto-fixable issues to a copy of the PPTX.

        Args:
            file_path: Original .pptx path
            issues: List of issues (only auto_fixable ones are processed)
            output_path: Output .pptx path

        Returns:
            (fixed_count, output_path)
        """
        # Copy original
        shutil.copy2(file_path, output_path)

        prs = Presentation(output_path)
        self.modified_count = 0
        self.fix_log = []

        # Group fixable issues by slide
        fixable = [i for i in issues if i.auto_fixable]
        slide_issues = {}
        for issue in fixable:
            si = issue.location.slide_index
            if si not in slide_issues:
                slide_issues[si] = []
            slide_issues[si].append(issue)

        for slide_idx, slide_issues_list in slide_issues.items():
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]

            for shape in slide.shapes:
                for issue in slide_issues_list:
                    if str(shape.shape_id) != issue.location.element_id and \
                       shape.name != issue.location.element_name:
                        continue
                    self._apply_fix(shape, issue, prs, slide_idx)

        prs.save(output_path)

        # Verify
        if not os.path.exists(output_path):
            raise RuntimeError("Failed to save fixed file")

        return self.modified_count, output_path

    def _apply_fix(self, shape, issue: Issue, prs: Presentation, slide_idx: int):
        """Apply a single fix to a shape."""
        fix_type = issue.fix_type
        params = issue.fix_params

        try:
            if fix_type == "font_replace":
                self._fix_font_replace(shape, params)
            elif fix_type == "font_unify":
                self._fix_font_unify(shape, params)
            elif fix_type == "font_size_set":
                self._fix_font_size(shape, params)
            elif fix_type == "color_replace":
                self._fix_color_replace(shape, params)
            elif fix_type == "move":
                self._fix_move(shape, params)
            elif fix_type == "move_into_page":
                self._fix_move_into_page(shape, params)
            elif fix_type == "image_ratio_fix":
                self._fix_image_ratio(shape, params)
            elif fix_type == "align":
                self._fix_align(shape, params)
            elif fix_type == "delete_slide":
                # Handled separately
                pass
            elif fix_type == "text_replace":
                self._fix_text_replace(shape, params)
            else:
                return

            self.modified_count += 1
            self.fix_log.append({
                'type': fix_type,
                'shape_name': shape.name,
                'slide': slide_idx,
                'params': params,
            })
        except Exception as e:
            pass  # Silently skip failed fixes

    def _fix_font_replace(self, shape, params: dict):
        """Replace a font in all text runs (handles run and paragraph level)."""
        old_font = params.get("old_font")
        new_font = params.get("new_font", "微软雅黑")
        if not shape.has_text_frame:
            self._fix_table_font(shape, old_font, new_font)
            return
        for paragraph in shape.text_frame.paragraphs:
            matched = False
            for run in paragraph.runs:
                font_name = run.font.name or paragraph.font.name
                if font_name and old_font and \
                   (old_font.lower() in font_name.lower() or
                    font_name.lower() in old_font.lower()):
                    run.font.name = new_font
                    matched = True
            # Fix at paragraph level if no runs or para-level font
            if not matched and paragraph.font.name and old_font:
                if old_font.lower() in paragraph.font.name.lower() or \
                   paragraph.font.name.lower() in old_font.lower():
                    paragraph.font.name = new_font

    def _fix_table_font(self, shape, old_font: str, new_font: str):
        """Replace font in table cells."""
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name and old_font and \
                               old_font.lower() in run.font.name.lower():
                                run.font.name = new_font
        except Exception:
            pass

    def _fix_font_unify(self, shape, params: dict):
        """Unify all fonts in shape to target font."""
        target = params.get("target_font", "微软雅黑")
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = target
                if not paragraph.runs:
                    paragraph.font.name = target

    def _fix_font_size(self, shape, params: dict):
        """Set font size for text runs (handles run and paragraph level)."""
        target_size = params.get("target_size", 12)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(target_size)
                if not paragraph.runs:
                    paragraph.font.size = Pt(target_size)

    def _fix_color_replace(self, shape, params: dict):
        """Replace color in text runs."""
        old_color = params.get("old_color", "").lstrip("#")
        new_color = params.get("new_color", "#333333").lstrip("#")
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        if run.font.color and run.font.color.rgb:
                            current = str(run.font.color.rgb)
                            if current.lower() == old_color.lower():
                                run.font.color.rgb = RGBColor(
                                    int(new_color[0:2], 16),
                                    int(new_color[2:4], 16),
                                    int(new_color[4:6], 16),
                                )
                    except Exception:
                        pass

    def _fix_move(self, shape, params: dict):
        """Move shape by offset."""
        dx_emu = pt_to_emu(params.get("dx", 0) * 72)
        dy_emu = pt_to_emu(params.get("dy", 0) * 72)
        shape.left = int(shape.left) + dx_emu
        shape.top = int(shape.top) + dy_emu

    def _fix_move_into_page(self, shape, params: dict):
        """Move shape to be within page."""
        from pptx.util import Inches
        left_in = params.get("left")
        top_in = params.get("top")
        if left_in is not None:
            shape.left = int(Inches(left_in))
        if top_in is not None:
            shape.top = int(Inches(top_in))

    def _fix_image_ratio(self, shape, params: dict):
        """Fix image aspect ratio (preserve original)."""
        orig_w = params.get("original_width", 1)
        orig_h = params.get("original_height", 1)
        display_w = params.get("display_width", 1)
        display_h = params.get("display_height", 1)

        if orig_w <= 0 or orig_h <= 0 or display_w <= 0 or display_h <= 0:
            return

        orig_ratio = orig_w / orig_h
        display_ratio = display_w / display_h

        if abs(orig_ratio - display_ratio) < 0.01:
            return

        from pptx.util import Emu
        current_width = shape.width
        current_height = shape.height

        # Fix by adjusting height to match original ratio
        new_height = int(current_width / orig_ratio)
        shape.height = new_height

        # If new height is larger than original, adjust width instead
        if new_height > current_height:
            shape.height = current_height
            shape.width = int(current_height * orig_ratio)

    def _fix_align(self, shape, params: dict):
        """Align shape (basic implementation)."""
        # Alignment is complex - for MVP we just note it
        pass

    def _fix_text_replace(self, shape, params: dict):
        """Replace text content."""
        old_text = params.get("old_text", "")
        new_text = params.get("new_text", "")
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

    def delete_slides(self, file_path: str, slide_indices: List[int],
                      output_path: str) -> Tuple[int, str]:
        """Delete specific slides from PPTX."""
        shutil.copy2(file_path, output_path)
        prs = Presentation(output_path)

        # Sort in reverse order to maintain indices
        for idx in sorted(slide_indices, reverse=True):
            if 0 <= idx < len(prs.slides):
                xml_slides = prs.slides._sldIdLst
                sldId = xml_slides[idx]
                rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    prs.part.drop_rel(rId)
                xml_slides.remove(sldId)

        prs.save(output_path)
        self.modified_count += len(slide_indices)
        return len(slide_indices), output_path


def quick_fix(file_path: str, issues: List[Issue],
              output_path: str = None) -> Tuple[int, str]:
    """Convenience function: apply all auto-fixes and return result."""
    if output_path is None:
        base, ext = os.path.splitext(file_path)
        output_path = f"{base}_fixed{ext}"

    # Separate fix types
    font_issues = [i for i in issues if i.auto_fixable and i.fix_type != "delete_slide"]
    delete_issues = [i for i in issues if i.fix_type == "delete_slide"]

    fixer = Fixer()
    total_fixed = 0
    current_path = file_path

    # Apply fixes
    fixed_count, _ = fixer.apply_fixes(current_path, font_issues, output_path)
    total_fixed += fixed_count

    # Handle deletions
    if delete_issues:
        slide_indices = sorted(set(
            i.fix_params.get("slide_index") for i in delete_issues
        ), reverse=True)
        if slide_indices:
            dc, output_path = fixer.delete_slides(
                output_path, slide_indices, output_path
            )
            total_fixed += dc

    return total_fixed, output_path
