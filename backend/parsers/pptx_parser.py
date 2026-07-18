"""PPTX parsing module - extracts all element data from .pptx files."""

from typing import Optional
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from backend.utils.geometry import emu_to_inches, emu_to_pt
from backend.utils.color import parse_pptx_color


@dataclass
class TextElement:
    """Parsed text element with formatting."""
    text: str
    font_name: Optional[str] = None
    font_size: Optional[float] = None  # in pt
    font_color: Optional[str] = None  # hex string
    bold: bool = False
    italic: bool = False
    is_title: bool = False
    paragraph_index: int = 0
    run_index: int = 0


@dataclass
class ShapeElement:
    """Parsed shape element with position and formatting."""
    shape_id: int
    name: str
    shape_type: str
    left: float  # inches
    top: float
    width: float
    height: float
    rotation: float = 0.0

    # Fill
    fill_color: Optional[str] = None
    fill_transparency: float = 0.0

    # Line/border
    line_color: Optional[str] = None
    line_width: Optional[float] = None

    # Text content
    text_elements: list = field(default_factory=list)
    has_text: bool = False
    text_overflow: bool = False

    # Group info
    group_id: Optional[str] = None


@dataclass
class ImageElement:
    """Parsed image element."""
    shape_id: int
    name: str
    left: float
    top: float
    width: float  # displayed width in inches
    height: float  # displayed height in inches
    image_width: int  # original width in px
    image_height: int  # original height in px
    content_type: Optional[str] = None
    blob_size: int = 0


@dataclass
class TableElement:
    """Parsed table element."""
    shape_id: int
    name: str
    left: float
    top: float
    width: float
    height: float
    rows: int = 0
    cols: int = 0


@dataclass
class SlideData:
    """All data extracted from a single slide."""
    slide_index: int
    width: float  # inches
    height: float
    shapes: list = field(default_factory=list)
    images: list = field(default_factory=list)
    tables: list = field(default_factory=list)
    all_text: str = ""
    background_color: Optional[str] = None
    layout_name: Optional[str] = None


@dataclass
class PptxData:
    """Complete extracted data from a PPTX file."""
    file_path: str
    file_size: int
    slide_width: float
    slide_height: float
    slides: list = field(default_factory=list)  # list of SlideData
    all_fonts: set = field(default_factory=set)
    all_colors: set = field(default_factory=set)
    all_images: list = field(default_factory=list)


def parse_pptx(file_path: str) -> Optional[PptxData]:
    """Parse a .pptx file and extract all relevant data."""
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Failed to open PPTX file: {e}")

    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)

    import os
    file_size = os.path.getsize(file_path)

    result = PptxData(
        file_path=file_path,
        file_size=file_size,
        slide_width=slide_width,
        slide_height=slide_height,
    )

    for idx, slide in enumerate(prs.slides):
        slide_data = SlideData(
            slide_index=idx,
            width=slide_width,
            height=slide_height,
            layout_name=slide.slide_layout.name if slide.slide_layout else None,
        )

        for shape in slide.shapes:
            # Parse shape element
            left = emu_to_inches(shape.left) if shape.left else 0
            top = emu_to_inches(shape.top) if shape.top else 0
            width = emu_to_inches(shape.width) if shape.width else 0
            height = emu_to_inches(shape.height) if shape.height else 0

            # Determine shape type
            stype = _get_shape_type(shape)

            if stype == "image" or stype == "picture":
                img = _parse_image(shape, left, top, width, height)
                if img:
                    slide_data.images.append(img)
                    slide_data.shapes.append(ShapeElement(
                        shape_id=shape.shape_id,
                        name=shape.name,
                        shape_type=stype,
                        left=left, top=top, width=width, height=height,
                    ))

            elif stype == "table":
                slide_data.tables.append(TableElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    left=left, top=top, width=width, height=height,
                    rows=shape.table.rows.__len__(),
                    cols=shape.table.columns.__len__(),
                ))
                slide_data.shapes.append(ShapeElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    shape_type="table",
                    left=left, top=top, width=width, height=height,
                ))
                _extract_table_text(shape, slide_data)

            elif shape.has_text_frame:
                texts = _extract_text_runs(shape)
                shape_el = ShapeElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    shape_type=stype,
                    left=left, top=top, width=width, height=height,
                    has_text=len(texts) > 0,
                    text_elements=texts,
                )

                # Check fill
                try:
                    if hasattr(shape, 'fill'):
                        fill_color = parse_pptx_color(shape.fill)
                        if fill_color:
                            shape_el.fill_color = fill_color
                            result.all_colors.add(fill_color)
                except Exception:
                    pass

                # Check text overflow
                shape_el.text_overflow = _check_text_overflow(shape)

                slide_data.shapes.append(shape_el)

                # Collect all text
                for t in texts:
                    slide_data.all_text += t.text + " "
                    if t.font_name:
                        result.all_fonts.add(t.font_name)
                    if t.font_color:
                        result.all_colors.add(t.font_color)

            else:
                # Other shape (group, connector, etc.)
                slide_data.shapes.append(ShapeElement(
                    shape_id=shape.shape_id,
                    name=shape.name,
                    shape_type=stype or "unknown",
                    left=left, top=top, width=width, height=height,
                ))

        result.slides.append(slide_data)

    return result


def _get_shape_type(shape) -> str:
    """Determine shape type string."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "auto_shape"
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return "freeform"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
            return "linked_picture"
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"
        elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            return "embedded_ole"
        elif shape.shape_type == MSO_SHAPE_TYPE.LINKED_OLE_OBJECT:
            return "linked_ole"
        else:
            return "unknown"
    except Exception:
        return "unknown"


def _parse_image(shape, left, top, width, height) -> Optional[ImageElement]:
    """Extract image data from a picture shape."""
    try:
        image = shape.image
        img_width_px = image.size[0] if hasattr(image, 'size') and image.size else 0
        img_height_px = image.size[1] if hasattr(image, 'size') and image.size else 0
        content_type = image.content_type if hasattr(image, 'content_type') else None
        blob_size = len(image.blob) if hasattr(image, 'blob') and image.blob else 0

        return ImageElement(
            shape_id=shape.shape_id,
            name=shape.name,
            left=left, top=top, width=width, height=height,
            image_width=img_width_px,
            image_height=img_height_px,
            content_type=content_type,
            blob_size=blob_size,
        )
    except Exception:
        return None


def _get_run_property(run, paragraph, prop_name):
    """Get a font property from run, falling back to paragraph default."""
    try:
        run_val = getattr(run.font, prop_name)
        if run_val is not None:
            return run_val
    except Exception:
        pass

    # Fall back to paragraph-level default
    try:
        para_val = getattr(paragraph.font, prop_name)
        if para_val is not None:
            return para_val
    except Exception:
        pass

    return None


def _extract_text_runs(shape) -> list:
    """Extract all text runs with formatting from a shape."""
    elements = []
    try:
        tf = shape.text_frame
        for pi, paragraph in enumerate(tf.paragraphs):
            for ri, run in enumerate(paragraph.runs):
                if not run.text.strip():
                    continue

                font_name = _get_run_property(run, paragraph, 'name')
                font_size_raw = _get_run_property(run, paragraph, 'size')
                font_size = font_size_raw.pt if font_size_raw else None
                font_color = None
                bold = False
                italic = False

                try:
                    fc = parse_pptx_color(run.font.color)
                    if not fc:
                        # Fall back to paragraph-level color
                        fc = parse_pptx_color(paragraph.font.color)
                    if fc:
                        font_color = fc
                except Exception:
                    pass

                try:
                    bold_val = _get_run_property(run, paragraph, 'bold')
                    bold = bool(bold_val) if bold_val is not None else False
                except Exception:
                    pass

                try:
                    italic_val = _get_run_property(run, paragraph, 'italic')
                    italic = bool(italic_val) if italic_val is not None else False
                except Exception:
                    pass

                is_title = _is_likely_title(paragraph, shape)

                elements.append(TextElement(
                    text=run.text,
                    font_name=font_name,
                    font_size=font_size,
                    font_color=font_color,
                    bold=bold,
                    italic=italic,
                    is_title=is_title,
                    paragraph_index=pi,
                    run_index=ri,
                ))

            # Handle empty runs (text set at paragraph level)
            if not paragraph.runs and paragraph.text.strip():
                font_name = None
                font_size = None
                try:
                    font_name = paragraph.font.name
                except Exception:
                    pass
                try:
                    if paragraph.font.size:
                        font_size = paragraph.font.size.pt
                except Exception:
                    pass

                elements.append(TextElement(
                    text=paragraph.text,
                    font_name=font_name,
                    font_size=font_size,
                    is_title=_is_likely_title(paragraph, shape),
                    paragraph_index=pi,
                    run_index=0,
                ))

    except Exception:
        pass

    return elements


def _extract_table_text(shape, slide_data):
    """Extract text from table cells."""
    try:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            slide_data.all_text += run.text + " "
                            try:
                                fn = run.font.name
                                if fn:
                                    pass  # fonts collected at outer level
                            except Exception:
                                pass
    except Exception:
        pass


def _check_text_overflow(shape) -> bool:
    """Check if text potentially overflows the shape."""
    try:
        tf = shape.text_frame
        if not tf.paragraphs:
            return False

        # Check if auto-size is not set or shape is too small
        text_length = sum(len(p.text) for p in tf.paragraphs)
        if text_length == 0:
            return False

        # Rough heuristic: very long text in small shape
        shape_area = (shape.width or 1) * (shape.height or 1)
        char_density = text_length / max(shape_area, 1)

        # If more than ~1 char per 10000 EMU², likely overflow
        return char_density > 1.0 / 10000
    except Exception:
        return False


def _is_likely_title(paragraph, shape) -> bool:
    """Heuristic: is this paragraph likely a title?"""
    # Check if it's a title placeholder
    try:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
            phf = shape.placeholder_format
            if phf.idx == 0:  # Title placeholder
                return True
    except Exception:
        pass

    # Check if it's the first paragraph and text is relatively short
    text = paragraph.text.strip()
    if not text:
        return False

    # If shape is near the top of the slide and text is short
    try:
        if shape.top and shape.top < 2000000 and len(text) < 100:  # top half
            return True
    except Exception:
        pass

    return False
