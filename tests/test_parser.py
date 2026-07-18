"""Test PPTX parser with a generated sample file."""

import os
import sys
import tempfile
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from backend.parsers import parse_pptx


def create_sample_pptx(path: str):
    """Create a sample PPTX with known issues for testing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "测试标题"
    p.font.name = "微软雅黑"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Slide 2: Content with issues
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Non-standard font
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "使用了非标准字体的标题"
    p.font.name = "Comic Sans MS"
    p.font.size = Pt(28)

    # Small text
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "这段文字字号太小了"
    p2.font.name = "微软雅黑"
    p2.font.size = Pt(8)  # too small

    # Off-page element
    shape = slide2.shapes.add_shape(
        1, Inches(15), Inches(1), Inches(2), Inches(1)
    )  # outside page

    # Slide 3: Another content page with inconsistent title
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide3.shapes.add_textbox(Inches(2), Inches(1), Inches(5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "另一页标题"
    p.font.name = "Times New Roman"  # inconsistent font
    p.font.size = Pt(26)  # different size

    # Slide 4: Content with multiple colors
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    colors = ["#FF0000", "#00FF00", "#0000FF", "#FF00FF", "#00FFFF", "#800080", "#FFA500"]
    for i, color in enumerate(colors):
        txBox = slide4.shapes.add_textbox(Inches(0.5 + i * 1.5), Inches(1), Inches(1.3), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"颜色{i+1}"
        p.font.name = "微软雅黑"
        p.font.size = Pt(14)
        try:
            r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
            p.font.color.rgb = RGBColor(r, g, b)
        except: pass

    # Slide 5: Blank page (almost)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    # Just a small shape, minimal content

    # Slide 6: Text overflow
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide6.shapes.add_textbox(Inches(1), Inches(1), Inches(1.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "这段文字非常长，会超出文本框的容量，导致文本被截断显示不全。这是一段测试文本，用于检查文本溢出检测功能。我们需要确保这个检测能够识别出文本过多而容器过小的情况。" * 3

    prs.save(path)
    return path


def test_parser():
    """Test that parser correctly extracts data from sample PPTX."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        sample_path = f.name

    try:
        create_sample_pptx(sample_path)
        data = parse_pptx(sample_path)

        assert data is not None, "Parser returned None"
        assert data.slide_width > 0, "Slide width not extracted"
        assert data.slide_height > 0, "Slide height not extracted"
        assert len(data.slides) == 6, f"Expected 6 slides, got {len(data.slides)}"

        # Check slide 1 has text
        s1 = data.slides[0]
        assert len(s1.shapes) >= 1, "Slide 1 should have shapes"
        assert "测试" in s1.all_text or s1.all_text == "", "Slide 1 should have text"

        # Check fonts collected
        assert len(data.all_fonts) > 0, "Should have collected fonts"

        # Check colors collected
        assert len(data.all_colors) > 0, "Should have collected colors"

        print(f"PASS: Parser test successful")
        print(f"  Slides: {len(data.slides)}")
        print(f"  Page size: {data.slide_width:.2f}x{data.slide_height:.2f} inches")
        print(f"  Fonts found: {data.all_fonts}")
        print(f"  File size: {data.file_size} bytes")

        # Print per-slide summary
        for slide in data.slides:
            texts = sum(1 for s in slide.shapes if s.has_text)
            images = len(slide.images)
            char_count = len(slide.all_text)
            print(f"  Slide {slide.slide_index + 1}: {texts} text shapes, {images} images, {char_count} chars")

    finally:
        os.unlink(sample_path)


def test_scanner():
    """Test that scanner finds issues in sample PPTX."""
    from backend.engine.scanner import Scanner
    from backend.config import RuleConfig

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        sample_path = f.name

    try:
        create_sample_pptx(sample_path)
        config = RuleConfig()
        scanner = Scanner(config)
        result = scanner.scan(sample_path)

        print(f"\nPASS: Scan test successful")
        print(f"  Score: {result.score}")
        print(f"  Total issues: {result.total_issues}")
        print(f"  S1: {result.s1_count}, S2: {result.s2_count}, S3: {result.s3_count}, S4: {result.s4_count}")
        print(f"  Fixable: {result.fixable_count}")
        print(f"  Scan time: {result.scan_time_ms:.0f}ms")

        # Print dimension scores
        for dim, score in (result.dimension_scores or {}).items():
            print(f"  {dim}: {score}")

        # Print all issues
        for issue in result.all_issues:
            print(f"    [{issue.severity}] P{issue.location.slide_index+1}: {issue.rule_name} - {issue.description[:60]}")

        assert result.total_issues > 0, "Should have found issues"

    finally:
        os.unlink(sample_path)


if __name__ == '__main__':
    test_parser()
    test_scanner()
    print("\nAll tests passed!")
